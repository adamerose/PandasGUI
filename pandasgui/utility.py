"""Reusable utility functions"""

import pandas as pd
import numpy as np


# %% Reshaping

def pivot(df, keys, categories, data):
    return pivoted_df


# %%
def save_figs_to_ppt(figs, filename):
    from pptx import Presentation

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    prs.save('test.pptx')


def flatten_multiindex(mi, sep=' - ', format=None):
    if type(mi) == pd.core.indexes.base.Index:
        return mi
    elif type(mi) == pd.core.indexes.multi.MultiIndex:
        # Flatten multi-index headers
        if format == None:
            # Flattern by putting sep between each header value
            flat_index = [sep.join(col).strip(sep) for col in mi.values]
        else:
            # Flatten according to the provided format string
            flat_index = []
            for tuple in mi.values:

                placeholders = []
                for name in mi.names:
                    if name is None:
                        name = ''
                    name = '{' + str(name) + '}'
                    placeholders.append(name)

                # Check if index segment contains each placeholder
                if all([item != '' for item in tuple]):
                    # Replace placeholders in format with corresponding values
                    flat_name = format
                    for i, val in enumerate(tuple):  # Iterates over the values in this index segment
                        flat_name = flat_name.replace(placeholders[i], val)
                else:
                    # If the segment doesn't contain all placeholders, just join them with sep instead
                    flat_name = sep.join(tuple).strip(sep)
                flat_index.append(flat_name)
    else:
        raise TypeError

    return flat_index


if __name__ == '__main__':
    #### EXAMPLES ####
    from pandasgui import show

    X = False  # I'm just using "if X:" instead of "if False:" so my IDE doesn't complain about unreachable code
    #### flatten_multiindex ####
    if 1:
        arrays = [['bar', 'bar', 'baz', 'baz', 'foo', 'foo', 'qux', 'qux'],
                  ['one', 'two', 'one', 'two', 'one', 'two', 'one', 'two']]
        tuples = list(zip(*arrays))
        index = pd.MultiIndex.from_tuples(tuples, names=['first', 'second'])
        s = pd.Series(np.random.randn(8), index=index)
        show(s, nonblocking=True)
        s.index = flatten_multiindex(s.index)
        show(s)

    #### pivot ####
    if X:
        pass
    if X:
        df = pd.read_csv('sample_data/pokemon.csv')
        keys = ['Generation']
        categories = ['Type 1', 'Type 2']
        data = {'Attack': ['min', 'max'], 'Defense': ['mean']}
        # data = {'Attack': ['mean']}

        grouped = df.groupby(keys + categories)
        aggregated = grouped.agg(data)
        aggregated.columns.names = ['AggColumn', 'AggFunc']
        pivoted_df = aggregated.unstack(categories)
