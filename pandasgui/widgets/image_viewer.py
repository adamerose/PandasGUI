"""Widgets for viewing figures"""

from PyQt5 import QtWidgets
import matplotlib
from matplotlib.backends.backend_qt4agg import FigureCanvasQTAgg
from matplotlib.backends.backend_qt4agg import NavigationToolbar2QT
from matplotlib.figure import Figure
import matplotlib.pyplot as plt
import random

matplotlib.use('Qt5Agg')

class TabbedFigureViewer(QtWidgets.QWidget):
    """
    Displays a list of matplotlib Figures in a tabbed viewer
    Args:
        figs (List): List of matplotlib Figures
    """

    def __init__(self, figs=[]):

        super().__init__()

        self.setWindowTitle("Tabbed Image Viewer")
        self.figs = figs
        layout = QtWidgets.QVBoxLayout()
        self.setLayout(layout)

        # Create tab widget
        self.tab_widget = QtWidgets.QTabWidget()
        self.tab_widget.setMovable(True)
        self.tab_widget.setTabsClosable(True)

        # Add figures to tab widget
        for fig in figs:
            self.add_tab(fig)

        # Just some button connected to `plot` method
        self.button = QtWidgets.QPushButton('To PowerPoint')
        self.button.clicked.connect(self.to_powerpoint)

        # Add widgets to layout
        layout.addWidget(self.tab_widget)
        layout.addWidget(self.button)

        self.show()

        # Center window on screen
        screen = QtWidgets.QDesktopWidget().screenGeometry()
        size = self.geometry()
        self.move(int((screen.width() - size.width()) / 2), int((screen.height() - size.height()) / 2))

    def add_tab(self, fig):
        """Add a new tab containing a FigureViewer showing fig"""
        image_viewer = FigureViewer(fig)
        # Adds them to the tab_view
        tab = QtWidgets.QWidget()
        tab_layout = QtWidgets.QVBoxLayout()
        tab.setLayout(tab_layout)
        tab_layout.addWidget(image_viewer)
        try:
            fig_title = fig.texts[0].get_text()
        except IndexError:
            fig_title = f"untitled"
        self.tab_widget.addTab(tab, fig_title)

    def to_powerpoint(self):
        # Write images to ppt
        import os
        filename = 'test.ppt'
        try:
            os.remove(filename)
        except OSError:
            pass
        for fig in self.figs:
            to_ppt_slide(fig, filename, append=True)
        os.system("start " + filename)
        win = TabbedFigureViewer(figs)


class FigureViewer(QtWidgets.QWidget):
    """
    Displays a matplotlib figure and toolbar in a widget

    Args:
        fig (Figure): matplotlib Figure to show
    """

    def __init__(self, fig=Figure()):
        super().__init__()

        self.canvas = FigureCanvasQTAgg(fig)
        self.toolbar = NavigationToolbar2QT(self.canvas, self)

        # Just some button connected to `plot` method
        self.button = QtWidgets.QPushButton('Plot Random')
        self.button.clicked.connect(self.setRandomImage)

        self.layout = QtWidgets.QVBoxLayout()
        self.setLayout(self.layout)
        self.layout.addWidget(self.toolbar)
        self.layout.addWidget(self.canvas)
        # self.layout.addWidget(self.button)

        self.toolbar.setStyleSheet('background: transparent;')

    def setFigure(self, fig):
        """Replace the current figure shown in the widget with fig"""
        plt.close(self.canvas.figure)

        new_canvas = FigureCanvasQTAgg(fig)
        self.layout.replaceWidget(self.canvas, new_canvas)

        new_toolbar = NavigationToolbar2QT(new_canvas, self)
        self.layout.replaceWidget(self.toolbar, new_toolbar)

        self.toolbar.setParent(None)
        self.toolbar = new_toolbar
        self.canvas = new_canvas

    def setRandomImage(self):
        """Plot some random data and show it in the widget"""
        data = [random.random() for i in range(10)]
        plt.figure()
        plt.plot(data, '*-')
        self.setFigure(plt.gcf())
        self.canvas.draw()


def to_ppt_slide(fig, file_path, append=False, padding=0.5):
    from io import StringIO, BytesIO
    import pptx
    from pptx import Presentation
    from pptx.util import Inches

    # Create in-memory image stream and save figure to it
    image_stream = BytesIO()
    fig.savefig(image_stream)

    if append:
        try:
            # Try opening the file if it already exists
            prs = Presentation(file_path)
        except pptx.exc.PackageNotFoundError:
            prs = Presentation()
    else:
        prs = Presentation()
    # Create a new slide with the blank template
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Center image without changing its aspect ratio
    slide_width = prs.slide_width.inches - 2 * padding
    slide_height = prs.slide_height.inches - 2 * padding
    fig_width, fig_height = fig.get_size_inches()
    if (fig_width / slide_width) > (fig_height / slide_height):
        # Image fits slide horizontally and must be scaled down vertically
        width = slide_width
        height = width * fig_height / fig_width
        top = padding + (slide_height - height) / 2
        left = padding
    else:
        # Image fits slide vertically and must be scaled down horizontally
        height = slide_height
        width = height * fig_width / fig_height
        left = padding + (slide_width - width) / 2
        top = padding

    # Convert from EMU to inches
    left = Inches(left)
    top = Inches(top)
    height = Inches(height)
    width = Inches(width)

    pic = slide.shapes.add_picture(image_stream, left, top, height=height, width=width)

    prs.save(file_path)


def ppt_section_slide(title, subtitle, file_path):
    from pptx import Presentation
    try:
        prs = Presentation(file_path)
    except:
        prs = Presentation(r'C:\Users\adrose\Desktop\AMD PowerPoint Template.pptx')
    picture_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(picture_slide_layout)

    # for x in slide.placeholders:
    #     print('%d %s' % (x.placeholder_format.idx, x.name))

    title_placeholder = slide.placeholders[0]
    subtitle_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    prs.save(file_path)


if __name__ == '__main__':
    from pandasgui.datasets import iris, flights, multi, pokemon

    app = QtWidgets.QApplication([])

    figs = []
    for df in [iris, flights, multi]:
        df.plot()
        figs.append(plt.gcf())

    win = FigureViewer(figs[0])
    win.show()

    win2 = TabbedFigureViewer(figs)
    win2.show()

    app.exec_()
